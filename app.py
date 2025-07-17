from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
import camelot
import pandas as pd
import uuid
import os
import logging
from werkzeug.utils import secure_filename
import pdf2image
import pytesseract
from PIL import Image
import io
import cv2
import numpy as np
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches
from datetime import datetime
import PyPDF2
import fitz  # PyMuPDF

app = Flask(__name__)
logging.basicConfig(level=logging.INFO)

CORS(app, origins=[
    "https://upload-agent.vercel.app",
    "https://upload-agent-git-main-ghalba-vieiras-projects-f2e9b128.vercel.app",
    "http://localhost:3000",
    "http://127.0.0.1:3000",
    "https://ghalba.app.n8n.cloud"
])

def extract_text_from_pdf(pdf_path):
    """Extrai texto completo do PDF com múltiplos métodos"""
    try:
        text = ""
        
        # Método 1: PyMuPDF (melhor para texto)
        try:
            doc = fitz.open(pdf_path)
            for page_num in range(doc.page_count):
                page = doc.load_page(page_num)
                text += page.get_text() + "\n\n"
            doc.close()
            
            if text.strip():
                logging.info(f"PyMuPDF extraiu {len(text)} caracteres")
                return text
        except Exception as e:
            logging.warning(f"PyMuPDF falhou: {e}")
        
        # Método 2: PyPDF2
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n\n"
            
            if text.strip():
                logging.info(f"PyPDF2 extraiu {len(text)} caracteres")
                return text
        except Exception as e:
            logging.warning(f"PyPDF2 falhou: {e}")
        
        # Método 3: OCR (para PDFs escaneados)
        try:
            logging.info("Tentando OCR para extração de texto...")
            images = pdf2image.convert_from_path(pdf_path, dpi=300)
            for i, image in enumerate(images):
                logging.info(f"Processando página {i+1} com OCR")
                
                # Converter PIL para OpenCV
                opencv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
                gray = cv2.cvtColor(opencv_image, cv2.COLOR_BGR2GRAY)
                
                # Aplicar filtros para melhorar OCR
                clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
                gray = clahe.apply(gray)
                
                # Threshold adaptativo
                thresh = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, 
                                             cv2.THRESH_BINARY, 11, 2)
                
                # OCR com configurações otimizadas
                custom_config = r'--oem 3 --psm 6 -l por'
                page_text = pytesseract.image_to_string(thresh, config=custom_config)
                text += page_text + "\n\n"
            
            logging.info(f"OCR extraiu {len(text)} caracteres")
            return text
        except Exception as e:
            logging.error(f"OCR falhou: {e}")
            return ""
    
    except Exception as e:
        logging.error(f"Erro na extração de texto: {e}")
        return ""

def extract_images_from_pdf(pdf_path):
    """Extrai e processa imagens do PDF"""
    try:
        doc = fitz.open(pdf_path)
        images_text = ""
        
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                try:
                    # Extrair imagem
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    if pix.n - pix.alpha < 4:  # GRAY ou RGB
                        # Converter para PIL Image
                        img_data = pix.tobytes("ppm")
                        pil_image = Image.open(io.BytesIO(img_data))
                        
                        # Converter para OpenCV
                        opencv_image = cv2.cvtColor(np.array(pil_image), cv2.COLOR_RGB2BGR)
                        gray = cv2.cvtColor(opencv_image, cv2.COLOR_BGR2GRAY)
                        
                        # Aplicar OCR na imagem
                        custom_config = r'--oem 3 --psm 6 -l por'
                        img_text = pytesseract.image_to_string(gray, config=custom_config)
                        
                        if img_text.strip():
                            images_text += f"\n[IMAGEM {page_num+1}-{img_index+1}]\n{img_text}\n"
                    
                    pix = None
                except Exception as e:
                    logging.warning(f"Erro ao processar imagem {img_index}: {e}")
                    continue
        
        doc.close()
        return images_text
    except Exception as e:
        logging.error(f"Erro na extração de imagens: {e}")
        return ""

def extract_tables_with_ocr(pdf_path):
    """Extrai tabelas usando OCR melhorado"""
    try:
        images = pdf2image.convert_from_path(pdf_path, dpi=300)
        all_data = []
        
        for page_num, image in enumerate(images):
            logging.info(f"Processando página {page_num + 1} com OCR para tabelas")
            
            # Converter para OpenCV
            opencv_image = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
            gray = cv2.cvtColor(opencv_image, cv2.COLOR_BGR2GRAY)
            
            # Detectar linhas horizontais e verticais (estrutura de tabela)
            horizontal_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (40, 1))
            vertical_kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (1, 40))
            
            # Aplicar filtros
            horizontal_lines = cv2.morphologyEx(gray, cv2.MORPH_OPEN, horizontal_kernel)
            vertical_lines = cv2.morphologyEx(gray, cv2.MORPH_OPEN, vertical_kernel)
            
            # Combinar linhas
            table_mask = cv2.addWeighted(horizontal_lines, 0.5, vertical_lines, 0.5, 0.0)
            
            # Se detectar estrutura de tabela, aplicar OCR otimizado
            if cv2.countNonZero(table_mask) > 100:
                # Configuração específica para tabelas
                custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyzÀÁÂÃÄÅÇÈÉÊËÌÍÎÏÑÒÓÔÕÖÙÚÛÜÝàáâãäåçèéêëìíîïñòóôõöùúûüý.,;:!?()[]{}\"\'\ -'
                text = pytesseract.image_to_string(gray, config=custom_config, lang='por')
                
                if text.strip():
                    lines = [line.strip() for line in text.split('\n') if line.strip()]
                    if lines:
                        page_data = {
                            'page': page_num + 1,
                            'text': text,
                            'lines': lines,
                            'has_table_structure': True
                        }
                        all_data.append(page_data)
        
        return all_data
    except Exception as e:
        logging.error(f"Erro no OCR de tabelas: {e}")
        return []

def convert_ocr_to_dataframe(ocr_data):
    """Converte dados OCR em DataFrame melhorado"""
    all_tables = []
    
    for page_data in ocr_data:
        lines = page_data['lines']
        
        # Identificar possíveis linhas de tabela
        table_lines = []
        for line in lines:
            # Separar por múltiplos espaços ou tabs
            parts = [part.strip() for part in line.split() if part.strip()]
            if len(parts) >= 2:  # Pelo menos 2 colunas
                table_lines.append(parts)
        
        if table_lines:
            # Determinar número de colunas baseado na linha com mais elementos
            max_cols = max(len(row) for row in table_lines)
            
            # Padronizar todas as linhas
            standardized_rows = []
            for row in table_lines:
                # Completar com strings vazias se necessário
                while len(row) < max_cols:
                    row.append('')
                standardized_rows.append(row[:max_cols])
            
            if standardized_rows:
                # Primeira linha como cabeçalho se parecer com cabeçalho
                if len(standardized_rows) > 1:
                    headers = standardized_rows[0]
                    data = standardized_rows[1:]
                else:
                    headers = [f'Coluna_{i+1}' for i in range(max_cols)]
                    data = standardized_rows
                
                df = pd.DataFrame(data, columns=headers)
                all_tables.append(df)
    
    return all_tables

def extract_tables_from_pdf(pdf_path):
    """Extrai tabelas do PDF com múltiplos métodos melhorados"""
    tables = []
    camelot_success = False
    
    # Método 1: Camelot Stream (melhor para tabelas sem bordas)
    try:
        tables_stream = camelot.read_pdf(pdf_path, pages="all", flavor="stream")
        if len(tables_stream) > 0:
            valid_tables = []
            for table in tables_stream:
                if not table.df.empty:
                    # Verificar se tem conteúdo significativo
                    non_empty_cells = table.df.astype(str).apply(
                        lambda x: x.str.strip().str.len() > 0
                    ).sum().sum()
                    
                    if non_empty_cells > 3:  # Pelo menos 3 células com conteúdo
                        valid_tables.append(table)
            
            if valid_tables:
                tables = valid_tables
                camelot_success = True
                logging.info(f"Camelot stream encontrou {len(tables)} tabelas válidas")
    except Exception as e:
        logging.warning(f"Camelot stream falhou: {e}")
    
    # Método 2: Camelot Lattice (melhor para tabelas com bordas)
    if not camelot_success:
        try:
            tables_lattice = camelot.read_pdf(pdf_path, pages="all", flavor="lattice")
            if len(tables_lattice) > 0:
                valid_tables = []
                for table in tables_lattice:
                    if not table.df.empty:
                        non_empty_cells = table.df.astype(str).apply(
                            lambda x: x.str.strip().str.len() > 0
                        ).sum().sum()
                        
                        if non_empty_cells > 3:
                            valid_tables.append(table)
                
                if valid_tables:
                    tables = valid_tables
                    camelot_success = True
                    logging.info(f"Camelot lattice encontrou {len(tables)} tabelas válidas")
        except Exception as e:
            logging.warning(f"Camelot lattice falhou: {e}")
    
    # Método 3: OCR melhorado para tabelas
    if not camelot_success:
        logging.info("Tentando extração de tabelas com OCR...")
        try:
            ocr_data = extract_tables_with_ocr(pdf_path)
            if ocr_data:
                ocr_tables = convert_ocr_to_dataframe(ocr_data)
                
                tables = []
                for df in ocr_tables:
                    if not df.empty:
                        class OCRTable:
                            def __init__(self, dataframe):
                                self.df = dataframe
                        
                        tables.append(OCRTable(df))
                
                logging.info(f"OCR encontrou {len(tables)} tabelas")
        except Exception as e:
            logging.error(f"OCR de tabelas falhou: {e}")
    
    return tables

@app.route("/convert/summary", methods=["POST"])
def generate_summary():
    """Extrai texto completo do PDF incluindo imagens"""
    pdf_path = None
    
    try:
        # Verificar se o arquivo foi enviado
        if "file" not in request.files:
            return jsonify({"error": "Nenhum arquivo PDF enviado"}), 400
        
        pdf_file = request.files["file"]
        
        if pdf_file.filename == '':
            return jsonify({"error": "Nenhum arquivo selecionado"}), 400
        
        # Salvar arquivo
        filename = secure_filename(pdf_file.filename)
        pdf_path = f"/tmp/{uuid.uuid4()}_{filename}"
        pdf_file.save(pdf_path)
        
        logging.info(f"Extraindo texto para resumo: {pdf_path}")
        
        # Extrair texto principal
        text = extract_text_from_pdf(pdf_path)
        
        # Extrair texto de imagens
        images_text = extract_images_from_pdf(pdf_path)
        
        # Combinar textos
        full_text = text + "\n\n" + images_text if images_text else text
        
        if not full_text.strip():
            return jsonify({
                "error": "Não foi possível extrair texto do PDF",
                "details": "O PDF pode estar protegido, corrompido ou não conter texto extraível"
            }), 422
        
        # Retornar texto extraído
        return jsonify({
            "text": full_text,
            "filename": filename,
            "extracted_at": datetime.now().isoformat(),
            "char_count": len(full_text),
            "has_images": len(images_text) > 0
        })
        
    except Exception as e:
        logging.error(f"Erro na extração de texto: {str(e)}")
        return jsonify({
            "error": "Erro interno do servidor",
            "details": str(e)
        }), 500
    
    finally:
        try:
            if pdf_path and os.path.exists(pdf_path):
                os.remove(pdf_path)
        except Exception as e:
            logging.warning(f"Erro na limpeza: {e}")

@app.route("/convert/excel", methods=["POST"])
def convert_to_excel():
    """Conversão melhorada para Excel"""
    pdf_path = None
    xlsx_path = None
    
    try:
        if "file" not in request.files:
            return jsonify({"error": "Nenhum arquivo PDF enviado"}), 400
        
        pdf_file = request.files["file"]
        
        if pdf_file.filename == '':
            return jsonify({"error": "Nenhum arquivo selecionado"}), 400
        
        filename = secure_filename(pdf_file.filename)
        pdf_path = f"/tmp/{uuid.uuid4()}_{filename}"
        pdf_file.save(pdf_path)
        
        logging.info(f"Convertendo para Excel: {pdf_path}")
        
        # Extrair tabelas
        tables = extract_tables_from_pdf(pdf_path)
        
        # Extrair texto das imagens para incluir no Excel
        images_text = extract_images_from_pdf(pdf_path)
        
        if not tables and not images_text:
            return jsonify({
                "error": "Nenhuma tabela ou conteúdo extraível encontrado",
                "details": "O PDF não contém tabelas ou imagens com texto extraível"
            }), 422
        
        xlsx_path = f"/tmp/{uuid.uuid4()}.xlsx"
        
        with pd.ExcelWriter(xlsx_path, engine='openpyxl') as writer:
            sheet_count = 0
            
            # Adicionar tabelas
            for i, table in enumerate(tables):
                df = table.df.copy()
                df = df.dropna(how='all', axis=0)  # Remover linhas completamente vazias
                df = df.dropna(how='all', axis=1)  # Remover colunas completamente vazias
                
                # Limpar dados
                df = df.applymap(lambda x: str(x).strip() if pd.notna(x) else '')
                
                # Remover linhas que são apenas espaços
                df = df[df.apply(lambda x: x.str.strip().str.len().sum() > 0, axis=1)]
                
                if not df.empty:
                    sheet_name = f"Tabela_{i+1}"[:31]
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
                    sheet_count += 1
                    logging.info(f"Tabela {i+1} salva: {len(df)} linhas x {len(df.columns)} colunas")
            
            # Adicionar texto de imagens se houver
            if images_text:
                lines = images_text.split('\n')
                image_data = []
                
                for line in lines:
                    if line.strip():
                        image_data.append([line.strip()])
                
                if image_data:
                    df_images = pd.DataFrame(image_data, columns=['Texto_Imagens'])
                    df_images.to_excel(writer, sheet_name='Texto_Imagens', index=False)
                    sheet_count += 1
            
            # Se não houver dados, criar uma planilha vazia com mensagem
            if sheet_count == 0:
                df_empty = pd.DataFrame([['Nenhuma tabela ou texto extraível encontrado']], 
                                      columns=['Resultado'])
                df_empty.to_excel(writer, sheet_name='Resultado', index=False)
        
        return send_file(
            xlsx_path,
            as_attachment=True,
            download_name=f"tabelas_{filename.replace('.pdf', '')}.xlsx",
            mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )
        
    except Exception as e:
        logging.error(f"Erro geral: {str(e)}")
        return jsonify({
            "error": "Erro interno do servidor",
            "details": str(e)
        }), 500
    
    finally:
        try:
            if pdf_path and os.path.exists(pdf_path):
                os.remove(pdf_path)
            if xlsx_path and os.path.exists(xlsx_path):
                os.remove(xlsx_path)
        except Exception as e:
            logging.warning(f"Erro na limpeza: {e}")

@app.route("/convert/word", methods=["POST"])
def convert_to_word():
    """Conversão melhorada para Word"""
    pdf_path = None
    docx_path = None
    
    try:
        if "file" not in request.files:
            return jsonify({"error": "Nenhum arquivo PDF enviado"}), 400
        
        pdf_file = request.files["file"]
        
        if pdf_file.filename == '':
            return jsonify({"error": "Nenhum arquivo selecionado"}), 400
        
        filename = secure_filename(pdf_file.filename)
        pdf_path = f"/tmp/{uuid.uuid4()}_{filename}"
        pdf_file.save(pdf_path)
        
        logging.info(f"Convertendo para Word: {pdf_path}")
        
        # Extrair texto principal
        text = extract_text_from_pdf(pdf_path)
        
        # Extrair texto de imagens
        images_text = extract_images_from_pdf(pdf_path)
        
        # Extrair tabelas
        tables = extract_tables_from_pdf(pdf_path)
        
        if not text.strip() and not images_text and not tables:
            return jsonify({
                "error": "Não foi possível extrair conteúdo do PDF",
                "details": "O PDF pode estar protegido ou não conter texto extraível"
            }), 422
        
        # Criar documento Word
        doc = Document()
        
        # Título
        title = doc.add_heading('Documento Convertido do PDF', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Metadados
        doc.add_paragraph(f"Arquivo original: {pdf_file.filename}")
        doc.add_paragraph(f"Convertido em: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}")
        doc.add_paragraph("")
        
        # Adicionar texto principal
        if text.strip():
            doc.add_heading('Conteúdo Principal', level=1)
            paragraphs = text.split('\n\n')
            for paragraph in paragraphs:
                if paragraph.strip():
                    cleaned_paragraph = paragraph.strip().replace('\n', ' ')
                    doc.add_paragraph(cleaned_paragraph)
        
        # Adicionar texto de imagens
        if images_text:
            doc.add_heading('Texto Extraído de Imagens', level=1)
            image_paragraphs = images_text.split('\n\n')
            for paragraph in image_paragraphs:
                if paragraph.strip():
                    cleaned_paragraph = paragraph.strip().replace('\n', ' ')
                    doc.add_paragraph(cleaned_paragraph)
        
        # Adicionar tabelas
        if tables:
            doc.add_heading('Tabelas Extraídas', level=1)
            
            for i, table in enumerate(tables):
                doc.add_heading(f'Tabela {i+1}', level=2)
                
                df = table.df.copy()
                df = df.dropna(how='all', axis=0)
                df = df.dropna(how='all', axis=1)
                df = df.fillna('')
                
                # Limpar dados
                df = df.applymap(lambda x: str(x).strip() if pd.notna(x) else '')
                
                if not df.empty:
                    try:
                        # Criar tabela no Word
                        word_table = doc.add_table(rows=len(df) + 1, cols=len(df.columns))
                        word_table.style = 'Table Grid'
                        
                        # Adicionar cabeçalhos
                        for j, column in enumerate(df.columns):
                            cell = word_table.cell(0, j)
                            cell.text = str(column) if column else f"Coluna {j+1}"
                            cell.paragraphs[0].runs[0].bold = True
                        
                        # Adicionar dados
                        for row_idx, (_, row) in enumerate(df.iterrows()):
                            for col_idx, value in enumerate(row):
                                cell = word_table.cell(row_idx + 1, col_idx)
                                cell.text = str(value) if pd.notna(value) else ""
                    except Exception as e:
                        logging.warning(f"Erro ao criar tabela {i+1}: {e}")
                        doc.add_paragraph(f"[Erro ao processar tabela {i+1}]")
                
                doc.add_paragraph("")
        
        # Salvar documento
        docx_path = f"/tmp/{uuid.uuid4()}.docx"
        doc.save(docx_path)
        
        return send_file(
            docx_path,
            as_attachment=True,
            download_name=f"documento_{filename.replace('.pdf', '')}.docx",
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        
    except Exception as e:
        logging.error(f"Erro na conversão para Word: {str(e)}")
        return jsonify({
            "error": "Erro interno do servidor",
            "details": str(e)
        }), 500
    
    finally:
        try:
            if pdf_path and os.path.exists(pdf_path):
                os.remove(pdf_path)
            if docx_path and os.path.exists(docx_path):
                os.remove(docx_path)
        except Exception as e:
            logging.warning(f"Erro na limpeza: {e}")

#power point precisa de atualizacao
@app.route("/convert/powerpoint", methods=["POST"])
def convert_to_powerpoint():
    """Conversão para PowerPoint"""
    pdf_path = None
    pptx_path = None
    
    try:
        if "file" not in request.files:
            return jsonify({"error": "Nenhum arquivo PDF enviado"}), 400
        
        pdf_file = request.files["file"]
        
        if pdf_file.filename == '':
            return jsonify({"error": "Nenhum arquivo selecionado"}), 400
        
        filename = secure_filename(pdf_file.filename)
        pdf_path = f"/tmp/{uuid.uuid4()}_{filename}"
        pdf_file.save(pdf_path)
        
        logging.info(f"Convertendo para PowerPoint: {pdf_path}")
        
        # Extrair texto do PDF
        text = extract_text_from_pdf(pdf_path)
        
        if not text.strip():
            return jsonify({
                "error": "Não foi possível extrair texto do PDF",
                "details": "O PDF pode estar protegido ou não conter texto extraível"
            }), 422
        
        # Extrair tabelas
        tables = extract_tables_from_pdf(pdf_path)
        
        # Criar apresentação PowerPoint
        prs = Presentation()
        
        # Slide de título
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Apresentação Convertida do PDF"
        subtitle.text = f"Arquivo: {pdf_file.filename}\nConvertido em: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}"
        
        # Dividir texto em slides
        if text.strip():
            paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
            
            # Agrupar parágrafos em slides (máximo 5 por slide)
            slide_groups = []
            current_group = []
            
            for paragraph in paragraphs[:20]:  # Limitar a 20 parágrafos
                if len(current_group) < 5:
                    current_group.append(paragraph)
                else:
                    slide_groups.append(current_group)
                    current_group = [paragraph]
            
            if current_group:
                slide_groups.append(current_group)
            
            # Criar slides de conteúdo
            for i, group in enumerate(slide_groups):
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                title.text = f"Conteúdo - Página {i+1}"
                
                # Adicionar texto ao slide
                text_frame = content.text_frame
                text_frame.clear()
                
                for j, paragraph in enumerate(group):
                    if j == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = paragraph[:200] + "..." if len(paragraph) > 200 else paragraph
                    p.level = 0
        
        # Adicionar tabelas em slides separados
        if tables:
            for i, table in enumerate(tables):
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                title.text = f"Tabela {i+1}"
                
                df = table.df.copy()
                df = df.dropna(how='all')
                df = df.fillna('')
                
                if not df.empty:
                    # Limitar tamanho da tabela para caber no slide
                    max_rows = min(10, len(df))
                    max_cols = min(6, len(df.columns))
                    
                    df_limited = df.iloc[:max_rows, :max_cols]
                    
                    # Adicionar tabela ao slide
                    left = PptxInches(1)
                    top = PptxInches(2)
                    width = PptxInches(8)
                    height = PptxInches(4)
                    
                    table_shape = slide.shapes.add_table(
                        rows=len(df_limited) + 1,
                        cols=len(df_limited.columns),
                        left=left,
                        top=top,
                        width=width,
                        height=height
                    )
                    
                    table_obj = table_shape.table
                    
                    # Adicionar cabeçalhos
                    for j, column in enumerate(df_limited.columns):
                        cell = table_obj.cell(0, j)
                        cell.text = str(column) if column else f"Col {j+1}"
                    
                    # Adicionar dados
                    for i, row in df_limited.iterrows():
                        for j, value in enumerate(row):
                            cell = table_obj.cell(i + 1, j)
                            cell.text = str(value)[:50] if pd.notna(value) else ""
        
        # Salvar apresentação
        pptx_path = f"/tmp/{uuid.uuid4()}.pptx"
        prs.save(pptx_path)
        
        return send_file(
            pptx_path,
            as_attachment=True,
            download_name="apresentacao_convertida.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except Exception as e:
        logging.error(f"Erro na conversão para PowerPoint: {str(e)}")
        return jsonify({
            "error": "Erro interno do servidor",
            "details": str(e)
        }), 500
    
    finally:
        try:
            if pdf_path and os.path.exists(pdf_path):
                os.remove(pdf_path)
            if pptx_path and os.path.exists(pptx_path):
                os.remove(pptx_path)
        except Exception as e:
            logging.warning(f"Erro na limpeza: {e}")

@app.route("/convert/summary", methods=["POST"])
def generate_summary():
    """Extrai texto do PDF para ser usado no frontend"""
    pdf_path = None
    
    try:
        if "file" not in request.files:
            return jsonify({"error": "Nenhum arquivo PDF enviado"}), 400
        
        pdf_file = request.files["file"]
        
        if pdf_file.filename == '':
            return jsonify({"error": "Nenhum arquivo selecionado"}), 400
        
        filename = secure_filename(pdf_file.filename)
        pdf_path = f"/tmp/{uuid.uuid4()}_{filename}"
        pdf_file.save(pdf_path)
        
        logging.info(f"Extraindo texto para resumo: {pdf_path}")
        
        # Extrair texto do PDF
        text = extract_text_from_pdf(pdf_path)
        
        if not text.strip():
            return jsonify({
                "error": "Não foi possível extrair texto do PDF",
                "details": "O PDF pode estar protegido ou não conter texto extraível"
            }), 422
        
        # Retornar texto extraído para o frontend processar
        return jsonify({
            "filename": secure_filename(pdf_file.filename),
            "summary": text or "Resumo não disponível.", 
            "extracted_at": datetime.now().isoformat()
        })
        
    except Exception as e:
        logging.error(f"Erro na extração de texto: {str(e)}")
        return jsonify({
            "error": "Erro interno do servidor",
            "details": str(e)
        }), 500
    
    finally:
        try:
            if pdf_path and os.path.exists(pdf_path):
                os.remove(pdf_path)
        except Exception as e:
            logging.warning(f"Erro na limpeza: {e}")

@app.route("/health", methods=["GET"])
def health_check():
    """Health check endpoint"""
    return jsonify({
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "services": {
            "pdf_processing": "active",
            "text_extraction": "active",
            "table_extraction": "active"
        }
    })

# Manter o endpoint original para compatibilidade
@app.route("/convert", methods=["POST"])
def convert_pdf():
    """Endpoint original mantido para compatibilidade"""
    return convert_to_excel()

if __name__ == "__main__":
    port = int(os.environ.get('PORT', 10000))
    app.run(host="0.0.0.0", port=port, debug=False)